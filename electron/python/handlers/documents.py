from __future__ import annotations

from pathlib import Path
from typing import Any, Callable

from .paths import default_class_folder, ensure_class_layout


def _docs_dir(params: dict[str, Any]) -> Path:
    if params.get("carpeta_clase"):
        return Path(params["carpeta_clase"]) / "03_documentos"
    if params.get("carpeta"):
        return Path(params["carpeta"])
    root = default_class_folder(str(params.get("tema") or "clase"), str(params.get("grado") or ""))
    ensure_class_layout(root)
    return root / "03_documentos"


def crear_pdf(params: dict[str, Any]) -> dict[str, Any]:
    titulo = str(params.get("titulo") or "Documento OnniVers")
    contenido = str(params.get("contenido") or params.get("texto") or "")
    out_dir = _docs_dir(params)
    out_dir.mkdir(parents=True, exist_ok=True)
    safe = "".join(c if c.isalnum() or c in "._- " else "_" for c in titulo)[:60].strip() or "documento"
    out = out_dir / f"{safe}.pdf"

    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas

        c = canvas.Canvas(str(out), pagesize=letter)
        width, height = letter
        y = height - 72
        c.setFont("Helvetica-Bold", 14)
        c.drawString(72, y, titulo[:90])
        y -= 36
        c.setFont("Helvetica", 11)
        for line in contenido.splitlines() or ["(sin contenido)"]:
            if y < 72:
                c.showPage()
                y = height - 72
                c.setFont("Helvetica", 11)
            c.drawString(72, y, line[:95])
            y -= 16
        c.save()
    except ImportError:
        out.write_text(f"{titulo}\n\n{contenido}", encoding="utf-8")
        out = out.with_suffix(".txt")
        return {
            "ok": True,
            "accion": "crear_pdf",
            "carpeta": str(out_dir),
            "archivos": [str(out)],
            "mensaje": "reportlab no instalado; se creó .txt. pip install reportlab",
        }

    return {
        "ok": True,
        "accion": "crear_pdf",
        "carpeta": str(out_dir),
        "archivos": [str(out)],
        "mensaje": "PDF creado",
    }


def crear_word(params: dict[str, Any]) -> dict[str, Any]:
    titulo = str(params.get("titulo") or "Documento")
    contenido = str(params.get("contenido") or "")
    out_dir = _docs_dir(params)
    out_dir.mkdir(parents=True, exist_ok=True)
    safe = "".join(c if c.isalnum() or c in "._- " else "_" for c in titulo)[:60].strip() or "documento"
    out = out_dir / f"{safe}.docx"
    try:
        from docx import Document

        doc = Document()
        doc.add_heading(titulo, level=1)
        for para in contenido.splitlines() or [""]:
            doc.add_paragraph(para)
        doc.save(str(out))
    except ImportError:
        out = out_dir / f"{safe}.txt"
        out.write_text(f"{titulo}\n\n{contenido}", encoding="utf-8")
        return {
            "ok": True,
            "accion": "crear_word",
            "archivos": [str(out)],
            "carpeta": str(out_dir),
            "mensaje": "python-docx no instalado; TXT creado. pip install python-docx",
        }
    return {"ok": True, "accion": "crear_word", "archivos": [str(out)], "carpeta": str(out_dir), "mensaje": "Word creado"}


def crear_excel(params: dict[str, Any]) -> dict[str, Any]:
    titulo = str(params.get("titulo") or "Hoja")
    filas = params.get("filas") or [["Pregunta", "Respuesta"], ["1", ""]]
    out_dir = _docs_dir(params)
    out_dir.mkdir(parents=True, exist_ok=True)
    safe = "".join(c if c.isalnum() or c in "._- " else "_" for c in titulo)[:60].strip() or "hoja"
    out = out_dir / f"{safe}.xlsx"
    try:
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = titulo[:31]
        for row in filas:
            ws.append(list(row))
        wb.save(str(out))
    except ImportError:
        out = out_dir / f"{safe}.csv"
        lines = [",".join(map(str, row)) for row in filas]
        out.write_text("\n".join(lines), encoding="utf-8")
        return {
            "ok": True,
            "accion": "crear_excel",
            "archivos": [str(out)],
            "carpeta": str(out_dir),
            "mensaje": "openpyxl no instalado; CSV creado. pip install openpyxl",
        }
    return {"ok": True, "accion": "crear_excel", "archivos": [str(out)], "carpeta": str(out_dir), "mensaje": "Excel creado"}


def crear_ppt(params: dict[str, Any]) -> dict[str, Any]:
    titulo = str(params.get("titulo") or params.get("tema") or "Clase")
    bullets = params.get("bullets") or [
        "Objetivos de la clase",
        "Contenido principal",
        "Actividad",
        "Evaluación",
        "Cierre",
    ]
    out_dir = _docs_dir(params)
    out_dir.mkdir(parents=True, exist_ok=True)
    safe = "".join(c if c.isalnum() or c in "._- " else "_" for c in titulo)[:60].strip() or "clase"
    out = out_dir / f"{safe}.pptx"
    try:
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = titulo
        body = prs.slides.add_slide(prs.slide_layouts[1])
        body.shapes.title.text = "Contenido"
        tf = body.shapes.placeholders[1].text_frame
        for bullet in bullets:
            tf.add_paragraph().text = str(bullet)
        prs.save(str(out))
    except ImportError:
        out = out_dir / f"{safe}_ppt.txt"
        out.write_text(titulo + "\n\n" + "\n".join(f"- {b}" for b in bullets), encoding="utf-8")
        return {
            "ok": True,
            "accion": "crear_ppt",
            "archivos": [str(out)],
            "carpeta": str(out_dir),
            "mensaje": "python-pptx no instalado; TXT creado. pip install python-pptx",
        }
    return {"ok": True, "accion": "crear_ppt", "archivos": [str(out)], "carpeta": str(out_dir), "mensaje": "PowerPoint creado"}


def leer_pdf(params: dict[str, Any]) -> dict[str, Any]:
    archivo = Path(params.get("archivo") or "")
    if not archivo.exists():
        return {"ok": False, "accion": "leer_pdf", "mensaje": "PDF no encontrado"}
    try:
        from pypdf import PdfReader

        reader = PdfReader(str(archivo))
        texto = "\n".join((page.extract_text() or "") for page in reader.pages)
    except ImportError:
        texto = archivo.read_bytes()[:2000].decode("utf-8", errors="ignore")
    except Exception as exc:  # noqa: BLE001
        return {"ok": False, "accion": "leer_pdf", "mensaje": str(exc)}
    return {
        "ok": True,
        "accion": "leer_pdf",
        "texto": texto[:12000],
        "archivos": [str(archivo)],
        "mensaje": f"{len(texto)} caracteres extraídos",
    }


def extraer_texto_pdf(params: dict[str, Any]) -> dict[str, Any]:
    return leer_pdf(params)


def convertir_a_pdf(params: dict[str, Any]) -> dict[str, Any]:
    contenido = str(params.get("contenido") or params.get("texto") or "")
    if params.get("archivo"):
        archivo = Path(params["archivo"])
        if archivo.suffix.lower() == ".txt":
            contenido = archivo.read_text(encoding="utf-8", errors="ignore")
    return crear_pdf({**params, "contenido": contenido, "titulo": params.get("titulo") or "Convertido"})


ACTIONS: dict[str, Callable[[dict[str, Any]], dict[str, Any]]] = {
    "crear_pdf": crear_pdf,
    "crear_word": crear_word,
    "crear_excel": crear_excel,
    "crear_ppt": crear_ppt,
    "leer_pdf": leer_pdf,
    "extraer_texto_pdf": extraer_texto_pdf,
    "convertir_a_pdf": convertir_a_pdf,
}
